from __future__ import annotations

import argparse
import hashlib
import json
import logging
import re
import sys
from pathlib import Path
from typing import Any, Iterable

from pydantic import BaseModel, Field, ValidationError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 兼容直接执行：
# python scripts/parse_ppt_kb.py --input ... --output ...
# PROJECT_ROOT = Path(__file__).resolve().parent.parent
from app.runtime import get_runtime_root

PROJECT_ROOT = get_runtime_root()

if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.logging_utils import setup_logging  # noqa: E402
from app.models import KBChunk  # noqa: E402

logger = logging.getLogger(__name__)

WHITESPACE_RE = re.compile(r"\s+")
QUESTION_LINE_RE = re.compile(r".{0,100}[?？].{0,20}$")
VALID_TEXT_CHAR_RE = re.compile(r"[\u4e00-\u9fffA-Za-z0-9]")
PUNCT_SPLIT_RE = re.compile(r"[。；;！!\n\r]+")

# 分类规则：按顺序匹配，前面的类别优先级更高
CATEGORY_RULES: list[tuple[str, tuple[str, ...]]] = [
    ("UKey", ("ukey", "u-key", "usbkey", "key", "uk", "证书介质")),
    ("证书", ("证书", "ca", "签名证书", "uk证书", "证书更新", "证书下载", "证书过期")),
    ("代发", ("代发", "工资", "批量代发", "代发工资")),
    ("权限", ("权限", "授权", "复核", "操作员", "审批", "管理员")),
    ("回单", ("回单", "电子回单", "回执", "回单打印", "回单下载")),
    ("登录", ("登录", "登陆", "登录失败", "验证码", "用户名", "密码", "认证")),
    ("控件", ("控件", "插件", "浏览器控件", "安全控件", "activex", "edge插件", "ie控件")),
    ("转账", ("转账", "付款", "汇款", "支付", "单笔转账", "批量转账")),
    ("账号权限", ("账号权限", "账户权限", "账户管理", "账号管理")),
    ("查询", ("查询", "明细", "余额", "流水", "对账")),
]


class ParsedSlideRecord(BaseModel):
    """
    单页 PPT 解析结果。

    chunk:
        标准 KBChunk 结构，可直接复用于后续索引、检索、问答流程。
    metadata:
        额外解析元信息。由于当前 KBChunk 模型未定义 metadata 字段，
        因此在导出 JSONL 时，metadata 会作为并列字段输出。
    """

    chunk: KBChunk
    metadata: dict[str, Any] = Field(default_factory=dict)


def normalize_text(value: Any) -> str:
    """
    文本归一化：
    1. None / 空值安全处理
    2. 全角空格转半角空格
    3. 连续空白压缩为单空格
    4. 去掉首尾空白
    """
    if value is None:
        return ""
    text = str(value).replace("\u3000", " ")
    text = WHITESPACE_RE.sub(" ", text).strip()
    return text


def iter_shapes_recursive(shapes: Iterable[Any]) -> Iterable[Any]:
    """
    递归遍历当前页中的全部 shape。
    用于兼容组图形（Group Shape）场景，避免漏掉组内文本或图片。
    """
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            subgroup = getattr(shape, "shapes", None)
            if subgroup is not None:
                yield from iter_shapes_recursive(subgroup)


def extract_text_from_shape(shape: Any) -> list[str]:
    """
    从单个 shape 中提取文本。

    支持：
    - 普通文本框 / 占位符
    - 表格单元格文本
    """
    results: list[str] = []

    # 普通文本框、占位符
    if getattr(shape, "has_text_frame", False):
        text = normalize_text(getattr(shape, "text", ""))
        if text:
            results.append(text)

    # 表格
    if getattr(shape, "has_table", False):
        table = getattr(shape, "table", None)
        if table is not None:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = normalize_text(cell.text)
                    if cell_text:
                        results.append(cell_text)

    return results


def ensure_image_dir(image_dir: Path) -> None:
    """
    确保图片输出目录存在。
    """
    try:
        image_dir.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        raise RuntimeError(f"创建图片目录失败：{image_dir}") from exc


def save_picture_shape(shape: Any, image_dir: Path, slide_no: int, picture_index: int) -> str | None:
    """
    保存单个图片 shape 到 data/parsed/images 目录。

    返回：
        相对项目根目录的路径字符串；若保存失败则返回 None。
    """
    image = getattr(shape, "image", None)
    if image is None:
        return None

    ext = normalize_text(getattr(image, "ext", "")).lower() or "png"
    if ext == "jpg":
        ext = "jpeg"

    file_name = f"slide_{slide_no:03d}_img_{picture_index:02d}.{ext}"
    output_path = image_dir / file_name

    try:
        output_path.write_bytes(image.blob)
    except OSError as exc:
        logger.warning("保存图片失败：slide=%s, file=%s, err=%s", slide_no, output_path, exc)
        return None

    try:
        relative_path = output_path.relative_to(PROJECT_ROOT)
        return relative_path.as_posix()
    except ValueError:
        return output_path.as_posix()


def extract_slide_images(slide: Any, image_dir: Path, slide_no: int) -> list[str]:
    """
    提取当前页内嵌图片并落盘。
    """
    image_paths: list[str] = []
    picture_index = 0

    for shape in iter_shapes_recursive(slide.shapes):
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            picture_index += 1
            saved = save_picture_shape(shape, image_dir=image_dir, slide_no=slide_no, picture_index=picture_index)
            if saved:
                image_paths.append(saved)

    return image_paths


def extract_slide_lines(slide: Any) -> list[str]:
    """
    提取单页全部文本，并尽量保持“按 shape 出现顺序”的线性文本结构。
    """
    lines: list[str] = []
    seen: set[str] = set()

    for shape in iter_shapes_recursive(slide.shapes):
        for text in extract_text_from_shape(shape):
            for part in re.split(r"[\r\n]+", text):
                clean = normalize_text(part)
                if not clean:
                    continue
                # 去重，避免某些 shape 同时被多处重复读取
                key = clean.casefold()
                if key in seen:
                    continue
                seen.add(key)
                lines.append(clean)

    return lines


def make_slide_text(lines: list[str]) -> str:
    """
    将单页全部文本整理为适合向量化与检索的自然语言串。
    """
    cleaned_lines = [normalize_text(line) for line in lines if normalize_text(line)]
    if not cleaned_lines:
        return ""
    return "。".join(cleaned_lines).strip("。") + "。"


def infer_title(lines: list[str], slide_no: int) -> str:
    """
    生成当前页标题。

    规则：
    1. 优先取首个问句
    2. 否则取首行有效文本
    3. 若没有文本，则退化为“PPT第X页”
    """
    valid_lines = [normalize_text(line) for line in lines if normalize_text(line)]
    if not valid_lines:
        return f"PPT第{slide_no}页"

    for line in valid_lines:
        if QUESTION_LINE_RE.search(line) or ("如何" in line) or ("怎么办" in line) or ("为什么" in line):
            return line[:40].rstrip()

    first_line = valid_lines[0]
    return first_line[:40].rstrip()


def classify_category(text: str) -> str | None:
    """
    基于关键词规则进行粗分类。
    """
    content = normalize_text(text).casefold()
    if not content:
        return None

    for category, keywords in CATEGORY_RULES:
        for keyword in keywords:
            if keyword.casefold() in content:
                return category

    return None


def count_effective_chars(text: str) -> int:
    """
    统计有效字符数量。
    用于判断该页是否“文本稀疏”。
    """
    return len(VALID_TEXT_CHAR_RE.findall(text))


def compute_chunk_hash(source_file: str, slide_no: int, slide_text: str) -> str:
    """
    为单页知识块生成稳定哈希。
    """
    raw = f"{source_file}||{slide_no}||{slide_text}"
    return hashlib.md5(raw.encode("utf-8")).hexdigest()


def build_full_text(title: str, category: str | None, slide_text: str) -> str:
    """
    构造适合 embedding 的自然语言文本。
    """
    parts: list[str] = [f"标题：{title}"]
    if category:
        parts.append(f"分类：{category}")
    if slide_text:
        parts.append(f"页面内容：{slide_text}")
    return "。".join(part.strip("。") for part in parts if part).strip("。") + "。"


def serialize_record(record: ParsedSlideRecord) -> dict[str, Any]:
    """
    序列化为可写入 JSONL 的字典结构。
    metadata 与 KBChunk 并列输出，便于后续调试与兼容。
    """
    payload = record.chunk.model_dump(mode="json", exclude_none=False)
    if record.metadata:
        payload["metadata"] = record.metadata
    return payload


def load_presentation(input_path: Path) -> Presentation:
    """
    读取 PPT 文件。
    """
    if not input_path.exists():
        raise FileNotFoundError(f"输入文件不存在：{input_path}")
    if not input_path.is_file():
        raise ValueError(f"输入路径不是文件：{input_path}")
    if input_path.suffix.lower() not in {".pptx", ".pptm"}:
        raise ValueError(f"暂不支持的 PPT 文件类型：{input_path.suffix}")

    try:
        return Presentation(str(input_path))
    except Exception as exc:
        raise RuntimeError(f"读取 PPT 失败：{input_path}") from exc


def parse_ppt_to_records(input_path: Path, image_dir: Path) -> list[ParsedSlideRecord]:
    """
    解析 PPT，并按“每页 slide 一个知识块”的原则生成记录。
    """
    ensure_image_dir(image_dir)
    prs = load_presentation(input_path)

    records: list[ParsedSlideRecord] = []
    skipped_empty = 0

    logger.info("开始解析 PPT：%s，页数=%s", input_path, len(prs.slides))

    for slide_index, slide in enumerate(prs.slides, start=1):
        try:
            lines = extract_slide_lines(slide)
            slide_text = make_slide_text(lines)
            image_paths = extract_slide_images(slide=slide, image_dir=image_dir, slide_no=slide_index)

            effective_chars = count_effective_chars(slide_text)
            text_sparse = effective_chars < 12

            # 完全空页：无文本且无图片，跳过
            if not slide_text and not image_paths:
                skipped_empty += 1
                logger.debug("跳过空白页：slide=%s", slide_index)
                continue

            title = infer_title(lines=lines, slide_no=slide_index)
            combined_for_classify = f"{title} {slide_text}"
            category = classify_category(combined_for_classify)
            full_text = build_full_text(title=title, category=category, slide_text=slide_text or "该页主要包含图片信息")
            chunk_hash = compute_chunk_hash(
                source_file=input_path.name,
                slide_no=slide_index,
                slide_text=slide_text or "|".join(image_paths),
            )

            chunk = KBChunk(
                doc_id=f"{input_path.stem}-slide-{slide_index}",
                source_file=input_path.name,
                source_type="ppt",
                title=title,
                category=category,
                question=None,
                answer=None,
                full_text=full_text,
                keywords=[kw for kw in [category, title] if kw],
                image_paths=image_paths,
                page_no=None,
                slide_no=slide_index,
                priority=0.95,
                chunk_hash=chunk_hash,
            )

            metadata: dict[str, Any] = {
                "text_sparse": text_sparse,
                "effective_char_count": effective_chars,
                "image_count": len(image_paths),
            }

            records.append(ParsedSlideRecord(chunk=chunk, metadata=metadata))

            logger.debug(
                "已解析 slide=%s, title=%s, category=%s, text_sparse=%s, images=%s",
                slide_index,
                title,
                category,
                text_sparse,
                len(image_paths),
            )

        except ValidationError as exc:
            logger.warning("Slide %s 的 KBChunk 校验失败，已跳过：%s", slide_index, exc)
        except Exception as exc:
            logger.exception("处理 Slide %s 时发生异常：%s", slide_index, exc)

    logger.info(
        "PPT 解析完成：总页数=%s，输出=%s，跳过空白页=%s，文件=%s",
        len(prs.slides),
        len(records),
        skipped_empty,
        input_path,
    )
    return records


def save_records_to_jsonl(records: list[ParsedSlideRecord], output_path: Path) -> None:
    """
    将解析结果保存为 JSONL。
    每行一个知识块对象；若存在 metadata，则一并输出。
    """
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        raise RuntimeError(f"创建输出目录失败：{output_path.parent}") from exc

    try:
        with output_path.open("w", encoding="utf-8") as f:
            for record in records:
                line = json.dumps(serialize_record(record), ensure_ascii=False)
                f.write(line + "\n")
    except OSError as exc:
        raise RuntimeError(f"写入 JSONL 失败：{output_path}") from exc

    logger.info("JSONL 已保存：%s（共 %s 条）", output_path, len(records))


def parse_args() -> argparse.Namespace:
    """
    解析命令行参数。
    """
    parser = argparse.ArgumentParser(
        description="解析企业网银问题带图.pptx，并按每页一个 KBChunk 输出为 JSONL。"
    )
    parser.add_argument(
        "--input",
        required=True,
        type=Path,
        help="输入 PPT 文件路径，例如 data/raw/企业网银问题带图.pptx",
    )
    parser.add_argument(
        "--output",
        required=True,
        type=Path,
        help="输出 JSONL 文件路径，例如 data/parsed/ppt_chunks.jsonl",
    )
    parser.add_argument(
        "--image-dir",
        type=Path,
        default=PROJECT_ROOT / "data" / "parsed" / "images",
        help="图片输出目录，默认 data/parsed/images/",
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        help="日志级别，例如 DEBUG / INFO / WARNING / ERROR",
    )
    return parser.parse_args()


def main() -> int:
    """
    命令行主入口。
    """
    args = parse_args()
    setup_logging(args.log_level, module_name=__name__)

    input_path: Path = args.input.expanduser().resolve()
    output_path: Path = args.output.expanduser().resolve()
    image_dir: Path = args.image_dir.expanduser().resolve()

    logger.info("参数：input=%s, output=%s, image_dir=%s", input_path, output_path, image_dir)

    try:
        records = parse_ppt_to_records(input_path=input_path, image_dir=image_dir)
        save_records_to_jsonl(records=records, output_path=output_path)
        return 0
    except Exception as exc:
        logger.exception("执行失败：%s", exc)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
